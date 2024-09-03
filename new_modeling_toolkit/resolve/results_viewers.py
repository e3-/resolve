import pathlib
from dataclasses import dataclass
from sys import platform
from typing import Optional

import pandas as pd
import pptx
import upath
import xlwings as xw
from loguru import logger
from PIL import ImageGrab

# Set traceback limit to 0 so that error message is more readable in Excel popup window

summary_sheet_mapping = [
    # Pull in some raw results as a band-aid
    # Further need to make sure that for any/all project specific RVs that come up all required input csvs that are needed are listed below

    ("constraints/ELCC_Facet_Constraint_LHS.csv", "ELCC_Facet_Constraint_LHS"),
    ("expressions/ELCC_Facet_Value.csv", "ELCC_Facet_Value"),
    ("variables/ELCC_MW.csv", "ELCC_MW"),
    ("results_summary/annual_load_components_summary.csv", "Load Components"),
    ("results_summary/asset_summary.csv", "Asset Summary"),
    ("constraints/Custom_Constraint.csv", "Custom Constraints"),
    ("results_summary/electrolyzer_summary.csv", "Electrolyzer"),
    ("results_summary/emissions_policy_resource_summary.csv", "Emissions Policy Resource"),
    ("results_summary/emissions_policy_tx_path_summary.csv", "Emissions Policy TxPath"),
    ("results_summary/energy_policy_summary.csv", "Energy Policy Summary"),
    ("results_summary/fuel_conversion_plant_summary.csv", "Fuel Conversion"),
    ("results_summary/fuel_storage_summary.csv", "Fuel Storage"),
    ("results_summary/fuel_transportation_summary.csv", "Fuel Transport"),
    ("results_summary/non_optimized_system_costs_summary.csv", "System"),
    ("results_summary/policy_summary.csv", "Policy Summary"),
    ("results_summary/prm_policy_resource_summary.csv", "PRM Resource Summary"),
    ("results_summary/resource_fuel_burn_summary.csv", "Resource Fuel Burn"),
    ("results_summary/resource_summary.csv", "Resource Summary"),
    ("results_summary/temporal_settings_summary.csv", "Temporal Settings"),
    ("results_summary/transmission_summary.csv", "Transmission Summary"),
    ("results_summary/zonal_summary.csv", "Zonal Summary"),
    ("results_summary/temporal_settings_summary.csv", "Temporal Settings"),
    ("temporal_settings/rep_periods.csv", "Representative Sample Days"),
    # ("temporal_settings/rep_periods.csv", "Representative Sample Days"),
    # ("resource_dispatch_summary.csv", "Dispatch"),
    # ("parameters/input_load_mw.csv", "Demand"),
]

def get_valid_results_folders(sheet_name: str, rng_name: str, wb: Optional["Book"] = None):
    """Get all RESOLVE results folders with non-empty ``results_summary`` subfolders."""

    if wb is None:
        wb = xw.Book.caller()

    results_path = upath.UPath(wb.fullname).parent / "reports" / "resolve"
    # Get all RESOLVE results folder names (nested list to make it "tall" instead of "wide"
    paths = [["/".join(p.parts[-3:-1])] for p in results_path.glob("**/results_summary") if any(p.iterdir())]
    # Write to Results Viewer range
    wb.sheets[sheet_name].range(rng_name).value = sorted(paths)


def populate_excel_results_viewer(case_name_and_timestamp: str, wb: Optional["Book"] = None):
    """Populate Excel Results Viewer tabs with summary data."""

    if wb is None:
        wb = xw.Book.caller()

    results_path = upath.UPath(case_name_and_timestamp)
    for summary_filename, target_sheet_name in summary_sheet_mapping:
        if (results_path / summary_filename).exists():
            if target_sheet_name not in [sheet.name for sheet in wb.sheets]:
                wb.sheets.add(target_sheet_name)
            wb.app.status_bar = f"Loading {summary_filename}"
            # Read summary CSV
            df = pd.read_csv(results_path / summary_filename)

            # Clear filters
            from sys import platform

            if platform == "win32":
                wb.sheets[target_sheet_name].api.AutoFilterMode = False
            elif platform == "darwin":
                wb.sheets[target_sheet_name].api.autofilter_mode = False

            # Clear data from target sheet
            wb.sheets[target_sheet_name].clear()

            # Write DataFrame to target sheet
            wb.sheets[target_sheet_name].range("A1").options(index=False, chunksize=100).value = df

    # Recalculate entire workbook
    wb.app.calculate()
    wb.app.status_bar = f"Finished loading {case_name_and_timestamp}"


def clear_data_sheets(wb: Optional["Book"] = None):
    """Clear raw CSV tabs in Results Viewer.

    TODO: Could streamline, since ``populate_excel_results_viewer`` also clears sheets.
    """

    if wb is None:
        wb = xw.Book.caller()

    for _, target_sheet_name in summary_sheet_mapping:
        if target_sheet_name not in [sheet.name for sheet in wb.sheets]:
            continue

        # Clear data from target sheet
        wb.sheets[target_sheet_name].api.AutoFilterMode = False
        wb.sheets[target_sheet_name].clear()

    wb.app.status_bar = "Data sheets cleared"


# TODO: Create an Enum for the slides & placeholders?


@dataclass
class ResultsAutomator:
    ppt: pptx.Presentation
    reports_folder: upath.UPath
    wb: xw.Book

    @property
    def case_name(self):
        return self.reports_folder.parts[-2]

    def save_figure_to_slide(self, *, sheet_name: str, figure_name: str, slide: pptx.slide.Slide, placeholder_idx: int):
        logger.debug(f"Creating slide for: {figure_name}")

        # Getting grouped API is different on Windows and macOS
        if platform == "win32":
            figure = self.wb.sheets[sheet_name].shapes[figure_name]
            figure.api.Copy()
        elif platform == "darwin":
            figure = self.wb.sheets[sheet_name].pictures[figure_name]
            figure.api.copy_object()

        img = ImageGrab.grabclipboard()
        # TODO: Seems like read/writing the figures to S3 isn't working. Need to build better workaround.
        img.save(f"{figure_name}.png")

        slide.placeholders[placeholder_idx].insert_picture(str(f"{figure.name}.png"))

        return slide

    def save_table_to_slide(self, *, sheet_name: str, table_name: str, slide: pptx.slide.Slide, placeholder_idx: int):
        xls_table = self.wb.sheets[sheet_name].range(table_name).options(pd.DataFrame, index=False, header=True).value

        xls_table = xls_table.loc[:, ~(xls_table == 0).all(axis=0)]
        xls_table = xls_table.loc[~(xls_table.iloc[:, 1:] == 0).all(axis=1), :]
        xls_table.columns = ["Resource Category"] + [col.strftime("%Y") for col in xls_table.columns[1:]]

        table = slide.placeholders[placeholder_idx].insert_table(xls_table.shape[0] + 1, xls_table.shape[1]).table
        table.last_row = True
        table.first_col = True

        for i, col in enumerate(table.columns):
            if i == 0:
                col.width = pptx.util.Inches(2.75)
            else:
                col.width = pptx.util.Inches(0.75)

        # Add the column names
        for i, column_name in enumerate(xls_table.columns):
            table.cell(0, i).text = column_name
            for p in table.cell(0, i).text_frame.paragraphs:
                p.font.size = pptx.util.Pt(11)
                p.alignment = pptx.enum.text.PP_ALIGN.CENTER

        # Add the data
        for i in range(xls_table.shape[0]):
            for j in range(xls_table.shape[1]):
                if j == 0:
                    table.cell(i + 1, j).text = str(xls_table.iloc[i, j])
                    for p in table.cell(i + 1, j).text_frame.paragraphs:
                        p.font.size = pptx.util.Pt(11)
                        p.alignment = pptx.enum.text.PP_ALIGN.LEFT
                else:
                    table.cell(i + 1, j).text = f"{xls_table.iloc[i, j]:.1f}"
                    for p in table.cell(i + 1, j).text_frame.paragraphs:
                        p.font.size = pptx.util.Pt(11)
                        p.alignment = pptx.enum.text.PP_ALIGN.CENTER

        table.height = pptx.util.Inches(3.6)

        return slide


def label_ppt_placeholders(template_path: pathlib.Path):
    """This cell creates a version of the CPUC Template PPT that labels each slide index and placeholder index
    (which is useful for adding new slides to the automation, because the pptx library operates only on the index numbers)

    """
    ppt = pptx.Presentation(template_path)
    for slide_layout in ppt.slide_layouts:
        slide = ppt.slides.add_slide(slide_layout)
        for shape in slide.placeholders:
            shape.text = f"{ppt.slide_layouts.index(slide_layout)} {slide_layout.name}: {shape.placeholder_format.idx} {shape.name}"

    ppt.save(template_path.with_stem(f"{template_path.stem} - placeholder names and index"))


# if __name__ == "__main__":
#     xw.Book(
#         "/Users/skramer/Library/CloudStorage/OneDrive-SharedLibraries-EnergyandEnvironmentalEconomics,Inc/CPUC IRP (1460) - Documents/RSP and PSP Analyses/RESOLVE 2023 PSP and 2024-25 TPP/Model UIs/Resolve Results Viewer - CPUC IRP 2023 PSP - PUBLIC - v1.xlsm"
#     ).set_mock_caller()
#     wb = xw.Book.caller()
#
#     populate_excel_results_viewer(
#         "reports/resolve/10_10_Hydrogen_InState_HourlyMatch_NoELCC_NoForce/2023-10-11 01-45-38/", wb=wb
#     )
