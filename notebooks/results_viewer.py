# ---
# jupyter:
#   jupytext:
#     custom_cell_magics: kql
#     text_representation:
#       extension: .py
#       format_name: percent
#       format_version: '1.3'
#       jupytext_version: 1.14.7
#   kernelspec:
#     display_name: Python 3 (ipykernel)
#     language: python
#     name: python3
# ---
# %% [markdown]
# # Results Viewer Controller
#
# This notebook serves two functions:
# 1. Create copies of the Resolve Results Template spreadsheet with case-specific results loaded.
# 2. Create a PowerPoint deck using the specified template to automatically generate simple chart & table slides.
# %%
from datetime import datetime

import ipywidgets as widgets
import pptx
import upath
import xlwings as xw
from loguru import logger

from resolve.resolve import results_viewers as rv

# %% [markdown]
# ### Controls & Settings
#
# #### _**<span style="color:red">Change these settings!</span>**_
#
# For your `results_path`, we recommend using the S3 path to your **run ID** (which is displayed after you use the `nmt submit` command). This will look like:
# - `s3://e3x-cpuc-irp-data/runs/[** your run ID **]/outputs/reports/resolve`

# %%
###################################
# UPDATE THESE SETTINGS EACH TIME #
###################################
# Template paths
sharepoint_path = upath.UPath(
    "/Users/rgo/Library/CloudStorage/OneDrive-SharedLibraries-EnergyandEnvironmentalEconomics,Inc/CPUC IRP (1460) - Documents/RSP and PSP Analyses/RESOLVE 2023 PSP and 2024-25 TPP/Model UIs"
)
ppt_template_path = sharepoint_path / "CPUC Template.pptx"
xlsx_template_path = sharepoint_path / "Resolve Results Viewer - CPUC IRP 2023 PSP - v0918.xlsx"

# Toggle whether to create PPT slides
create_ppt = True

# Define where the results are stored (by default, they are in AWS S3)
results_path = upath.UPath("s3://e3x-cpuc-irp-data/runs/rgo-20230904.1/outputs/reports/resolve")

###################################

if isinstance(results_path, upath.implementations.cloud.S3Path):
    logger.info("Your results seem to be in an S3 bucket. You will now be prompted to sign into AWS!")
    # !aws sso login


# %%
cases_list = ["/".join(p.parts[-3:-1]) for p in results_path.glob("**/results_summary") if any(p.iterdir())]

cases = widgets.SelectMultiple(
    options=sorted(cases_list),
    rows=10,
    description="Select Cases to Load: ",
    layout={
        "width": "max-content",
    },
    style={"description_width": "160px"},
)
display(cases)

# %%
# This list controls the slides to make, and the figures/tables to insert into each slide

slides_to_make = [
    dict(
        title="Planned & Selected Capacity (GW)",
        subtitle_index=14,
        slide_index=6,
        figures_and_tables=[
            dict(sheet_name="Build", figure_name="new_capacity_split", placeholder_idx=13),
        ],
    ),
    dict(
        title="Planned & Selected Capacity (GW)",
        subtitle_index=14,
        slide_index=12,
        figures_and_tables=[
            dict(sheet_name="Build", table_name="selected_capacity_table", placeholder_idx=15),
        ],
    ),
    dict(
        title="Planned & Selected Capacity, Compared to LSE Plans (GW)",
        subtitle_index=14,
        slide_index=6,
        figures_and_tables=[
            dict(sheet_name="Comparisons", figure_name="Group 4", placeholder_idx=13),
        ],
    ),
    dict(
        title="PRM Results",
        subtitle_index=15,
        slide_index=3,
        figures_and_tables=[
            dict(sheet_name="Reliability", figure_name="mtr_results", placeholder_idx=13),
            dict(sheet_name="Reliability", figure_name="prm_results", placeholder_idx=14),
        ],
    ),
    dict(
        title="RPS & SB 100",
        subtitle_index=15,
        slide_index=3,
        figures_and_tables=[
            dict(sheet_name="RPS & SB 100", figure_name="rps", placeholder_idx=13),
            dict(sheet_name="RPS & SB 100", figure_name="Group 6", placeholder_idx=14),
        ],
    ),
    dict(
        title="In-state & Unspecified Import Emissions (MMT)",
        subtitle_index=15,
        slide_index=3,
        figures_and_tables=[
            dict(sheet_name="Emissions", figure_name="ghg_emissions", placeholder_idx=14),
        ],
    ),
    dict(
        title="Selected Transmission Upgrades",
        subtitle_index=14,
        slide_index=2,
        figures_and_tables=[
            dict(sheet_name="Transmission", figure_name="tx_upgrades", placeholder_idx=13),
        ],
    ),
]

# %%
ppt = pptx.Presentation(upath.UPath(ppt_template_path))
upath.UPath("../reports/resolve/").mkdir(exist_ok=True, parents=True)

counter = 1

for _, case_path in enumerate(cases.value):
    if not upath.UPath(results_path / case_path / "results_summary" / "resource_summary.csv").exists():
        continue

    case_name = case_path.split("/")[0]
    with xw.App(visible=False):
        results_viewer = xw.Book(upath.UPath(xlsx_template_path))

        #################################
        # Load case into Results Viewer #
        #################################

        logger.info(f"Loading case {counter} of {len(cases.value)}: {case_name} from {case_path}")
        counter += 1

        reports_folder = upath.UPath(results_path / case_path)
        results_viewer.sheets["Dashboard"].range("active_case").value = str(reports_folder)
        rv.populate_excel_results_viewer(case_name_and_timestamp=reports_folder, wb=results_viewer)

        # Save case-specific Results Viewer to reports folder
        logger.info(f"Saving loaded results viewer")
        results_viewer.save(path=f"../reports/resolve/Resolve Results Viewer - {reports_folder.parts[-2]}.xlsx")

        # Create ResultsAutomator instance
        r = rv.ResultsAutomator(
            ppt=ppt,
            reports_folder=reports_folder,
            wb=results_viewer,
        )

        ###################
        ## Create slides ##
        ###################
        if create_ppt:
            ppt.slides[0].placeholders[1].text = f"\n{datetime.now().strftime('%B %d, %Y')}"

            # Section header
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            slide.shapes.title.text = f"{r.case_name}"

            for slide_config in slides_to_make:
                slide = ppt.slides.add_slide(ppt.slide_layouts[slide_config["slide_index"]])
                slide.placeholders[slide_config["subtitle_index"]].text = r.case_name
                slide.shapes.title.text = slide_config["title"]

                for fig_config in slide_config["figures_and_tables"]:
                    if "figure_name" in fig_config.keys():
                        slide = r.save_figure_to_slide(
                            sheet_name=fig_config["sheet_name"],
                            figure_name=fig_config["figure_name"],
                            slide=slide,
                            placeholder_idx=fig_config["placeholder_idx"],
                        )
                    elif "table_name" in fig_config.keys():
                        slide = r.save_table_to_slide(
                            sheet_name=fig_config["sheet_name"],
                            table_name=fig_config["table_name"],
                            slide=slide,
                            placeholder_idx=fig_config["placeholder_idx"],
                        )

            # Save PPT just because we can
            ppt.save(upath.UPath(f"../reports/resolve/{datetime.now().strftime('%Y-%m-%d')} Results Summary.pptx"))

        # Close spreadsheet
        results_viewer.close()

logger.info("Done!")

# %%
